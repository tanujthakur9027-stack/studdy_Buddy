"""
Document ingestion service
==========================
Responsibilities
- Parse PDF (pdfplumber → PyPDF2 fallback), TXT/MD, DOCX/DOC, PPT/PPTX, XLSX/XLS, images (OCR).
- Generate a short auto-description from the first ~500 chars of extracted text.
- Chunk text with LangChain's RecursiveCharacterTextSplitter.
- Maintain TWO vector stores per document:
    1. ChromaDB  – persisted to disk, survives restarts, used for cross-session search.
    2. FAISS     – in-process memory, zero-latency for the current session, per-doc index.
- Expose a unified `retrieve_context()` that merges and deduplicates results from both.
"""
from __future__ import annotations

import asyncio
import io
import logging
import re
import textwrap
import uuid
from pathlib import Path
from typing import Optional

import aiofiles
import pdfplumber
import PyPDF2
from langchain_core.documents import Document
from langchain_community.document_loaders import Docx2txtLoader
from langchain_community.vectorstores import FAISS
from langchain_chroma import Chroma
from langchain_community.embeddings import FastEmbedEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import get_settings
from utils.text_utils import clean_text, count_tokens

logger = logging.getLogger(__name__)
settings = get_settings()

# ── Singleton embeddings (fastembed ONNX — no PyTorch, no API key required) ───
_embeddings: Optional[FastEmbedEmbeddings] = None


def get_embeddings() -> FastEmbedEmbeddings:
    global _embeddings
    if _embeddings is None:
        _embeddings = FastEmbedEmbeddings(model_name="BAAI/bge-small-en-v1.5")
    return _embeddings


# ── In-process FAISS registry  ────────────────────────────────────────────────
_faiss_registry: dict[str, FAISS] = {}
_faiss_global: Optional[FAISS] = None

def _get_faiss_for_doc(doc_id: str) -> Optional[FAISS]:
    return _faiss_registry.get(doc_id)


def _get_faiss_global() -> Optional[FAISS]:
    return _faiss_global


# ── ChromaDB (persisted) ─────────────────────────────────────────────────────
def get_chroma(collection: str = "studybuddy") -> Chroma:
    return Chroma(
        collection_name=collection,
        embedding_function=get_embeddings(),
        persist_directory=settings.chroma_persist_dir,
    )


# ── PDF parsing ───────────────────────────────────────────────────────────────

def _extract_pdf_pdfplumber(file_bytes: bytes) -> list[tuple[int, str]]:
    pages: list[tuple[int, str]] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            raw = page.extract_text(x_tolerance=2, y_tolerance=2) or ""
            tables = page.extract_tables() or []
            table_text = ""
            for table in tables:
                for row in table:
                    row_cells = [str(c or "").strip() for c in row]
                    table_text += "  |  ".join(row_cells) + "\n"
            combined = raw + ("\n" + table_text if table_text.strip() else "")
            pages.append((i, clean_text(combined)))
    return pages


def _extract_pdf_pypdf2(file_bytes: bytes) -> list[tuple[int, str]]:
    pages: list[tuple[int, str]] = []
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    for i, page in enumerate(reader.pages, start=1):
        raw = page.extract_text() or ""
        pages.append((i, clean_text(raw)))
    return pages


def extract_pdf_pages(file_bytes: bytes, filename: str) -> tuple[list[tuple[int, str]], str]:
    try:
        pages = _extract_pdf_pdfplumber(file_bytes)
        total_chars = sum(len(t) for _, t in pages)
        if total_chars > 50:
            logger.info("[%s] pdfplumber extracted %d pages, %d chars", filename, len(pages), total_chars)
            return pages, "pdfplumber"
        logger.warning("[%s] pdfplumber gave sparse output (%d chars) — trying PyPDF2", filename, total_chars)
    except Exception as exc:
        logger.warning("[%s] pdfplumber failed (%s) — falling back to PyPDF2", filename, exc)

    pages = _extract_pdf_pypdf2(file_bytes)
    total_chars = sum(len(t) for _, t in pages)
    logger.info("[%s] PyPDF2 extracted %d pages, %d chars", filename, len(pages), total_chars)
    return pages, "PyPDF2"


# ── Plain-text / Markdown parsing ─────────────────────────────────────────────

def _strip_md_frontmatter(text: str) -> str:
    return re.sub(r"^---\s*\n.*?\n---\s*\n", "", text, flags=re.DOTALL)


def extract_txt(file_bytes: bytes, is_markdown: bool = False) -> list[tuple[int, str]]:
    raw = file_bytes.decode("utf-8", errors="replace")
    if is_markdown:
        raw = _strip_md_frontmatter(raw)
    segments = [s.strip() for s in re.split(r"\n{2,}", raw) if s.strip()]
    pages: list[tuple[int, str]] = []
    page_num = 1
    current: list[str] = []
    current_len = 0
    for seg in segments:
        current.append(seg)
        current_len += len(seg)
        if current_len >= 1500:
            pages.append((page_num, clean_text("\n\n".join(current))))
            page_num += 1
            current, current_len = [], 0
    if current:
        pages.append((page_num, clean_text("\n\n".join(current))))
    return pages


# ── PowerPoint parsing ────────────────────────────────────────────────────────

def extract_pptx(file_bytes: bytes) -> list[tuple[int, str]]:
    """Extract text from each slide of a PPT/PPTX file."""
    from pptx import Presentation  # lazy import

    prs = Presentation(io.BytesIO(file_bytes))
    pages: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            # Also pull table cell text
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append("  |  ".join(cells))
        text = clean_text("\n".join(parts))
        if text:
            pages.append((i, text))
    return pages


# ── Excel parsing ─────────────────────────────────────────────────────────────

def extract_xlsx(file_bytes: bytes) -> list[tuple[int, str]]:
    """Extract text from all sheets of an XLSX/XLS workbook."""
    import openpyxl  # lazy import

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    pages: list[tuple[int, str]] = []
    for sheet_num, sheet in enumerate(wb.worksheets, start=1):
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append("  |  ".join(cells))
        text = clean_text("\n".join(rows))
        if text:
            pages.append((sheet_num, f"[Sheet: {sheet.title}]\n{text}"))
    wb.close()
    return pages


# ── Image OCR parsing ─────────────────────────────────────────────────────────

def extract_image(file_bytes: bytes) -> list[tuple[int, str]]:
    """Run Tesseract OCR on an image file and return extracted text as a single page."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        raw = pytesseract.image_to_string(img)
        text = clean_text(raw)
        if text:
            return [(1, text)]
    except Exception as exc:
        logger.warning("OCR extraction failed: %s", exc)
    return []


# ── Description generator ─────────────────────────────────────────────────────

def generate_description(pages: list[tuple[int, str]], filename: str, parser_used: str) -> str:
    """
    Build a short human-readable description of the document from the first chunk of text.
    Format: '<N> pages/slides/rows · <parser> · <first 120 chars of content>…'
    """
    total = len(pages)
    # pick correct unit label
    ext = Path(filename).suffix.lower()
    if ext in (".ppt", ".pptx"):
        unit = "slide" if total == 1 else "slides"
    elif ext in (".xlsx", ".xls"):
        unit = "sheet" if total == 1 else "sheets"
    elif ext in (".png", ".jpg", ".jpeg", ".webp"):
        unit = "image"
    else:
        unit = "page" if total == 1 else "pages"

    # grab the first ~500 chars of combined text
    combined = " ".join(t for _, t in pages[:3])
    snippet = textwrap.shorten(combined, width=120, placeholder="…")

    return f"{total} {unit} · parsed via {parser_used} · {snippet}"


# ── Chunking ──────────────────────────────────────────────────────────────────

def make_splitter() -> RecursiveCharacterTextSplitter:
    return RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ". ", "! ", "? ", " ", ""],
    )


def pages_to_documents(
    pages: list[tuple[int, str]],
    doc_id: str,
    filename: str,
) -> list[Document]:
    docs: list[Document] = []
    for page_num, text in pages:
        if not text.strip():
            continue
        docs.append(Document(
            page_content=text,
            metadata={
                "doc_id": doc_id,
                "filename": filename,
                "page": page_num,
                "source": f"{filename}#page{page_num}",
            },
        ))
    return docs


def chunk_documents(docs: list[Document]) -> list[Document]:
    splitter = make_splitter()
    chunks = splitter.split_documents(docs)
    for i, chunk in enumerate(chunks):
        chunk.metadata["chunk_index"] = i
        chunk.metadata["token_count"] = count_tokens(chunk.page_content)
    return chunks


# ── Ingestion pipeline ────────────────────────────────────────────────────────

async def save_upload(file_bytes: bytes, filename: str) -> tuple[str, str]:
    """Persist raw file to disk. Returns (filepath, doc_id)."""
    upload_path = Path(settings.upload_dir)
    upload_path.mkdir(parents=True, exist_ok=True)
    doc_id = str(uuid.uuid4())
    ext = Path(filename).suffix
    filepath = upload_path / f"{doc_id}{ext}"
    async with aiofiles.open(filepath, "wb") as f:
        await f.write(file_bytes)
    return str(filepath), doc_id


class IngestionStats:
    """Carries per-document ingestion metrics back to the caller."""
    def __init__(
        self,
        doc_id: str,
        filename: str,
        pages: int,
        chunks: int,
        total_chars: int,
        total_tokens: int,
        parser_used: str,
        description: str = "",
    ):
        self.doc_id = doc_id
        self.filename = filename
        self.pages = pages
        self.chunks = chunks
        self.total_chars = total_chars
        self.total_tokens = total_tokens
        self.parser_used = parser_used
        self.description = description


def _sync_process_and_index(
    file_bytes: bytes,
    filepath: str,
    doc_id: str,
    filename: str,
) -> IngestionStats:
    """
    Synchronous ingestion pipeline — runs in a thread via asyncio.to_thread()
    so it never blocks the async event loop.

    1. Extract text based on file extension.
    2. Generate a short description from the extracted text.
    3. Convert to LangChain Documents, chunk, tag metadata.
    4. Index into ChromaDB (persistent) AND FAISS (in-memory, per-doc + global).
    Returns IngestionStats for the response body.
    """
    global _faiss_global

    ext = Path(filename).suffix.lower()
    parser_used = "direct"

    # ── 1. Extract ────────────────────────────────────────────────────────────
    if ext == ".pdf":
        pages, parser_used = extract_pdf_pages(file_bytes, filename)
    elif ext in (".txt", ".bin"):
        pages = extract_txt(file_bytes, is_markdown=False)
        parser_used = "plaintext"
    elif ext == ".md":
        pages = extract_txt(file_bytes, is_markdown=True)
        parser_used = "markdown"
    elif ext in (".docx", ".doc"):
        loader = Docx2txtLoader(filepath)
        loaded = loader.load()
        pages = [(i + 1, clean_text(doc.page_content)) for i, doc in enumerate(loaded)]
        parser_used = "python-docx"
    elif ext in (".pptx", ".ppt"):
        pages = extract_pptx(file_bytes)
        parser_used = "python-pptx"
    elif ext in (".xlsx", ".xls"):
        pages = extract_xlsx(file_bytes)
        parser_used = "openpyxl"
    elif ext in (".png", ".jpg", ".jpeg", ".webp"):
        pages = extract_image(file_bytes)
        parser_used = "tesseract-ocr"
    else:
        raise ValueError(f"Unsupported file extension: {ext}")

    if not pages or all(not t.strip() for _, t in pages):
        raise ValueError("No extractable text found in the uploaded file.")

    # ── 2. Auto-description ───────────────────────────────────────────────────
    description = generate_description(pages, filename, parser_used)

    # ── 3. Build Document objects & chunk ─────────────────────────────────────
    raw_docs = pages_to_documents(pages, doc_id, filename)
    chunks = chunk_documents(raw_docs)

    total_chars = sum(len(c.page_content) for c in chunks)
    total_tokens = sum(c.metadata.get("token_count", 0) for c in chunks)

    # ── 4a. ChromaDB (disk-persisted) ─────────────────────────────────────────
    Path(settings.chroma_persist_dir).mkdir(parents=True, exist_ok=True)
    chroma = get_chroma()
    chroma.add_documents(chunks)

    # ── 4b. FAISS per-document in-memory index ────────────────────────────────
    embeddings = get_embeddings()
    doc_faiss = FAISS.from_documents(chunks, embeddings)
    _faiss_registry[doc_id] = doc_faiss
    logger.info("[%s] FAISS per-doc index built (%d chunks)", filename, len(chunks))

    # ── 4c. Merge into global FAISS index ─────────────────────────────────────
    if _faiss_global is None:
        _faiss_global = FAISS.from_documents(chunks, embeddings)
    else:
        _faiss_global.merge_from(doc_faiss)
    logger.info("Global FAISS index now contains %d vectors", _faiss_global.index.ntotal)

    return IngestionStats(
        doc_id=doc_id,
        filename=filename,
        pages=len(pages),
        chunks=len(chunks),
        total_chars=total_chars,
        total_tokens=total_tokens,
        parser_used=parser_used,
        description=description,
    )


async def process_and_index(
    file_bytes: bytes,
    filepath: str,
    doc_id: str,
    filename: str,
) -> IngestionStats:
    """
    Async wrapper — offloads all blocking CPU/IO work to a thread pool
    so the FastAPI event loop stays responsive during ingestion.
    """
    return await asyncio.to_thread(
        _sync_process_and_index, file_bytes, filepath, doc_id, filename
    )


# ── Retrieval ─────────────────────────────────────────────────────────────────

def retrieve_context(
    query: str,
    doc_id: Optional[str] = None,
    k: int = 5,
) -> list[Document]:
    """
    Retrieve the top-k most relevant chunks for *query*.

    Strategy (in order of preference):
    1. If doc_id is given AND a per-doc FAISS index exists → use it (fastest, most precise).
    2. If doc_id is given but FAISS index is absent → fall back to global FAISS with metadata filter.
    3. If no doc_id → query global FAISS index.
    4. Final safety net → ChromaDB with optional metadata filter.
    """
    results: list[Document] = []

    if doc_id:
        per_doc = _get_faiss_for_doc(doc_id)
        if per_doc is not None:
            try:
                results = per_doc.similarity_search(query, k=k)
                logger.debug("FAISS per-doc hit: %d results for doc %s", len(results), doc_id)
            except Exception as exc:
                logger.warning("FAISS per-doc search failed: %s", exc)

    if not results:
        global_idx = _get_faiss_global()
        if global_idx is not None:
            try:
                candidates = global_idx.similarity_search(query, k=k * 3)
                if doc_id:
                    candidates = [d for d in candidates if d.metadata.get("doc_id") == doc_id]
                results = candidates[:k]
                logger.debug("FAISS global hit: %d results", len(results))
            except Exception as exc:
                logger.warning("FAISS global search failed: %s", exc)

    if not results:
        try:
            chroma = get_chroma()
            filter_dict = {"doc_id": doc_id} if doc_id else None
            results = chroma.similarity_search(query, k=k, filter=filter_dict)
            logger.debug("ChromaDB fallback: %d results", len(results))
        except Exception as exc:
            logger.error("ChromaDB search also failed: %s", exc)

    seen: set[int] = set()
    unique: list[Document] = []
    for doc in results:
        h = hash(doc.page_content)
        if h not in seen:
            seen.add(h)
            unique.append(doc)

    return unique[:k]


def list_indexed_docs() -> list[dict]:
    """Return metadata for all currently indexed documents (from FAISS registry)."""
    result = []
    for doc_id, idx in _faiss_registry.items():
        try:
            vectors = idx.index.ntotal
        except Exception:
            vectors = 0
        result.append({"doc_id": doc_id, "vectors": vectors})
    return result


def populate_faiss_from_chroma() -> int:
    """
    Rebuild all in-process FAISS indexes from ChromaDB on startup.

    Called once during the FastAPI lifespan so that vector search works
    at full speed from the very first request — even after a worker restart
    or a Streamlit Cloud rerun — without requiring re-uploads.

    Returns the total number of vectors loaded.
    """
    global _faiss_global

    try:
        chroma = get_chroma()
        # Fetch all stored documents from ChromaDB
        collection = chroma._collection
        raw = collection.get(include=["documents", "metadatas", "embeddings"])
    except Exception as exc:
        logger.warning("populate_faiss_from_chroma: ChromaDB unavailable — %s", exc)
        return 0

    ids        = raw.get("ids", [])
    texts      = raw.get("documents", [])
    metadatas  = raw.get("metadatas", []) or [{}] * len(ids)
    embeddings_raw = raw.get("embeddings")  # list of vectors or None

    if not ids:
        logger.info("populate_faiss_from_chroma: ChromaDB is empty — nothing to rebuild")
        return 0

    # Group by doc_id
    from collections import defaultdict
    groups: dict[str, list[tuple[str, str, dict]]] = defaultdict(list)
    for i, (text, meta) in enumerate(zip(texts, metadatas)):
        doc_id = (meta or {}).get("doc_id", "__unknown__")
        groups[doc_id].append((ids[i], text, meta or {}))

    emb_model  = get_embeddings()
    total      = 0

    for doc_id, items in groups.items():
        try:
            docs = [
                Document(page_content=text, metadata=meta)
                for _, text, meta in items
            ]
            doc_faiss = FAISS.from_documents(docs, emb_model)
            _faiss_registry[doc_id] = doc_faiss
            total += doc_faiss.index.ntotal

            # Merge into global index
            if _faiss_global is None:
                _faiss_global = FAISS.from_documents(docs, emb_model)
            else:
                _faiss_global.merge_from(doc_faiss)
        except Exception as exc:
            logger.warning("populate_faiss_from_chroma: failed to rebuild doc %s — %s", doc_id, exc)

    logger.info(
        "populate_faiss_from_chroma: rebuilt %d doc indexes, %d total vectors",
        len(groups), total,
    )
    return total
